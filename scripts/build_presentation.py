"""
Build the 14-slide capstone presentation for the ResQ HP Disaster Relief Agent.

Generates docs/ResQ_Capstone_Presentation.pptx — architecture, LangGraph nodes,
resource-allocation engine, human-in-the-loop email workflow, and the
"why this tech" rationale, all drawn from the actual code in src/resq_project.

Usage:
    python scripts/build_presentation.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ── Palette ────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x10, 0x2A, 0x43)
DARKNAVY  = RGBColor(0x0A, 0x1C, 0x30)
ORANGE    = RGBColor(0xF2, 0x54, 0x2D)
AMBER     = RGBColor(0xF5, 0xA6, 0x23)
TEAL      = RGBColor(0x16, 0x8A, 0x7A)
SKY       = RGBColor(0x2E, 0x6F, 0xA5)
LIGHT     = RGBColor(0xF4, 0xF6, 0xF9)
CARD      = RGBColor(0xEA, 0xEF, 0xF5)
INK       = RGBColor(0x22, 0x2E, 0x3A)
GRAY      = RGBColor(0x5A, 0x68, 0x74)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
RED       = RGBColor(0xC0, 0x3A, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.12):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    return sp


def set_text(shape_or_tf, runs, size=12, color=INK, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.0):
    """runs: str, or list of paragraphs; each paragraph is str or list of
    (text, {bold/color/size/italic}) run tuples."""
    tf = shape_or_tf.text_frame if hasattr(shape_or_tf, "text_frame") else shape_or_tf
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.09)
    tf.margin_right = Inches(0.09)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    if isinstance(runs, str):
        runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        if isinstance(para, str):
            para = [(para, {})]
        for text, style in para:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = FONT
            f.size = Pt(style.get("size", size))
            f.bold = style.get("bold", bold)
            f.italic = style.get("italic", False)
            f.color.rgb = style.get("color", color)
    return tf


def card(slide, x, y, w, h, title, body, fill=CARD, title_color=NAVY,
         body_size=10.5, title_size=12, accent=ORANGE):
    """Rounded card with a colored left keyline, bold title and bullet body."""
    sp = rect(slide, x, y, w, h, fill=fill)
    bar = rect(slide, x, y + Inches(0.08), Inches(0.055), h - Inches(0.16),
               fill=accent, shape=MSO_SHAPE.RECTANGLE)
    bar.left = x + Inches(0.05)
    paras = [[(title, {"bold": True, "size": title_size, "color": title_color})]]
    for line in body:
        if isinstance(line, str):
            paras.append([("▪ ", {"color": accent, "bold": True, "size": body_size}),
                          (line, {"size": body_size, "color": INK})])
        else:  # pre-styled runs
            paras.append(line)
    set_text(sp, paras, anchor=MSO_ANCHOR.TOP, space_after=3)
    return sp


def arrow(slide, x1, y1, x2, y2, color=GRAY, width=2.0, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.shadow.inherit = False
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if dashed:
        dash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(dash)
    head = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(head)
    return conn


def header(slide, num, title, subtitle=None):
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT, shape=MSO_SHAPE.RECTANGLE)
    band = rect(slide, 0, 0, SLIDE_W, Inches(1.02), fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    chip = rect(slide, Inches(0.35), Inches(0.24), Inches(0.55), Inches(0.55), fill=ORANGE)
    set_text(chip, [[(f"{num:02d}", {"bold": True, "size": 18, "color": WHITE})]],
             align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(Inches(1.05), Inches(0.10), Inches(11.6), Inches(0.9))
    paras = [[(title, {"bold": True, "size": 24, "color": WHITE})]]
    if subtitle:
        paras.append([(subtitle, {"size": 12.5, "color": RGBColor(0xBF, 0xD3, 0xE6)})])
    set_text(tb, paras, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    # footer
    ft = slide.shapes.add_textbox(Inches(0.35), Inches(7.12), Inches(12.6), Inches(0.3))
    set_text(ft, [[("ResQ · HP Disaster Relief Resource Matching Agent", {"size": 9, "color": GRAY}),
                   (f"   ·   {num} / 14", {"size": 9, "color": GRAY, "bold": True})]],
             anchor=MSO_ANCHOR.MIDDLE)
    return slide


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def flow_box(slide, x, y, w, h, title, sub, fill=NAVY, sub_color=None, title_size=11.5,
             sub_size=8.5):
    sp = rect(slide, x, y, w, h, fill=fill)
    sub_color = sub_color or RGBColor(0xC9, 0xD8, 0xE8)
    paras = [[(title, {"bold": True, "size": title_size, "color": WHITE})]]
    if sub:
        paras.append([(sub, {"size": sub_size, "color": sub_color})])
    set_text(sp, paras, align=PP_ALIGN.CENTER, space_after=1)
    return sp


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=DARKNAVY, shape=MSO_SHAPE.RECTANGLE)
rect(s, 0, Inches(5.0), SLIDE_W, Inches(0.06), fill=ORANGE, shape=MSO_SHAPE.RECTANGLE)
tb = s.shapes.add_textbox(Inches(0.9), Inches(1.35), Inches(11.5), Inches(2.6))
set_text(tb, [
    [("ResQ", {"bold": True, "size": 60, "color": ORANGE}),
     ("  —  AI-Assisted Disaster Relief Coordination", {"bold": True, "size": 34, "color": WHITE})],
    [("A LangGraph multi-agent system for Himachal Pradesh: one incident report in — "
      "an enriched, ranked, explainable, human-approved response plan out.",
      {"size": 17, "color": RGBColor(0xBF, 0xD3, 0xE6)})],
], anchor=MSO_ANCHOR.TOP, space_after=14)
chips = [
    ("6-agent LangGraph pipeline", ORANGE),
    ("RAG over 6 official HP datasets", TEAL),
    ("Explainable urgency scoring", SKY),
    ("Human-approved email dispatch", AMBER),
]
cx = Inches(0.9)
for label, col in chips:
    w = Inches(2.85)
    c = rect(s, cx, Inches(4.15), w, Inches(0.5), fill=col)
    set_text(c, [[(label, {"bold": True, "size": 11.5, "color": WHITE})]], align=PP_ALIGN.CENTER)
    cx += w + Inches(0.18)
tb = s.shapes.add_textbox(Inches(0.9), Inches(5.35), Inches(11.5), Inches(1.4))
set_text(tb, [
    [("Major Capstone Project · Applied AI", {"bold": True, "size": 15, "color": WHITE})],
    [("Team ResQ  ·  Python · Streamlit · LangGraph · LangChain · Ollama · ChromaDB",
      {"size": 12, "color": RGBColor(0x8F, 0xA6, 0xBC)})],
], anchor=MSO_ANCHOR.TOP, space_after=6)
notes(s, "Hook: in a disaster, minutes matter but information is scattered. ResQ turns a "
         "single incident report into a coordination-ready plan — with a human always in "
         "control of the final action.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM & IMPACT
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, 2, "The Problem: Disaster Triage Is an Information Problem",
       "Himachal Pradesh — flash floods, landslides, cloudbursts, GLOF, wildfire")
card(s, Inches(0.35), Inches(1.3), Inches(6.25), Inches(2.75),
     "What responders face", [
    "Affected people need Medical, Shelter, Food, Water, Rescue, Transport — all at once, from one location report.",
    "Kangra alone logged 4,027 landslides in 2023; Mandi 2,169 (HIMCOSTE inventory).",
    "Key NH corridors (NH-154 @ Kotrupi, NH-5 @ Nigulsari…) close chronically every monsoon — routes can't be assumed.",
    "Hospitals, NULM shelters, schools, CWC river stations and volunteers live in disconnected PDFs, CSVs and spreadsheets.",
], accent=RED)
card(s, Inches(6.85), Inches(1.3), Inches(6.15), Inches(2.75),
     "What control rooms actually need", [
    "A ranked shortlist of actionable resources — not a raw database dump.",
    "Context: live weather alert, district risk tier, GLOF & wildfire proneness.",
    "A route estimate with known road-risk warnings.",
    "An explainable urgency score to prioritise between incidents.",
    "An auditable, human-approved way to actually contact providers.",
], accent=TEAL)
band = rect(s, Inches(0.35), Inches(4.3), Inches(12.65), Inches(2.55), fill=NAVY)
set_text(band, [
    [("Our answer", {"bold": True, "size": 15, "color": ORANGE})],
    [("ResQ ingests one structured incident report (district, place, disaster type, needs), enriches it with live weather "
      "and risk intelligence, retrieves and ranks matching resources, plans a route, computes an explainable 0–100 urgency "
      "score, and drafts a coordination email — which a human reviews, edits and approves before anything is sent.",
      {"size": 13.5, "color": WHITE})],
    [("Design stance: decision support, never an autonomous dispatcher. Every output says “verify with local authorities”; "
      "every send is gated by a person and logged to an audit trail.",
      {"size": 12, "color": RGBColor(0xBF, 0xD3, 0xE6), "italic": True})],
], anchor=MSO_ANCHOR.MIDDLE, space_after=8)
notes(s, "Frame the problem as information fragmentation, not lack of resources. The numbers "
         "(4,027 landslides in Kangra) come from the HIMCOSTE 2023 inventory bundled in the repo.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, 3, "System Architecture — Four Layers, One Audit Trail",
       "Streamlit UI → LangGraph agents → retrieval & live context → human approval gate")
Y1, Y2, Y3, Y4 = Inches(1.28), Inches(2.62), Inches(3.96), Inches(5.42)
BH = Inches(1.12)
# Layer 1 — UI
flow_box(s, Inches(0.35), Y1, Inches(4.05), BH, "Streamlit App (app/app.py)",
         "Incident form · geocoded location · folium maps", fill=SKY)
flow_box(s, Inches(4.55), Y1, Inches(4.05), BH, "Volunteer Worklist",
         "needs.csv ↔ resources.csv match review", fill=SKY)
flow_box(s, Inches(8.75), Y1, Inches(4.25), BH, "Grounded RAG Chatbot",
         "answers only from ingested HP data", fill=SKY)
# Layer 2 — intelligence
flow_box(s, Inches(0.35), Y2, Inches(4.05), BH, "LangGraph Pipeline (workflow.py)",
         "6 agent nodes · shared DisasterState · conditional edges", fill=NAVY)
flow_box(s, Inches(4.55), Y2, Inches(4.05), BH, "Deterministic Allocation Engine",
         "coordination.py · 100-pt transparent scoring, no LLM", fill=NAVY)
flow_box(s, Inches(8.75), Y2, Inches(4.25), BH, "Relevance-Gated RAG (chatbot.py)",
         "refuse before LLM if best hit > 0.70 distance", fill=NAVY)
# Layer 3 — data & models
flow_box(s, Inches(0.35), Y3, Inches(4.05), BH, "ChromaDB (6 collections)",
         "hospitals · shelters · schools · CWC · knowledge · glacial lakes", fill=TEAL)
flow_box(s, Inches(4.55), Y3, Inches(4.05), BH, "Live Context APIs",
         "Open-Meteo weather · OpenRouteService routes · Nominatim geocoding", fill=TEAL)
flow_box(s, Inches(8.75), Y3, Inches(4.25), BH, "Local LLM — Ollama llama3.2:1b",
         "temperature 0.1 · ranking + report writing only", fill=TEAL)
# Layer 4 — human gate
g1 = flow_box(s, Inches(0.35), Y4, Inches(5.3), BH, "Human Approval Gate",
              "review → edit → approve / reject — nothing dispatched autonomously", fill=ORANGE)
g2 = flow_box(s, Inches(6.05), Y4, Inches(3.4), BH, "Gmail SMTP (smtplib)",
              "smtp.gmail.com:587 · STARTTLS · app password", fill=ORANGE)
g3 = flow_box(s, Inches(9.85), Y4, Inches(3.15), BH, "Audit Trail",
              "logs/approvals.jsonl · UTC-stamped decisions", fill=ORANGE)
for x in (Inches(2.35), Inches(6.55), Inches(10.85)):
    arrow(s, x, Y1 + BH, x, Y2, color=GRAY)
    arrow(s, x, Y2 + BH, x, Y3, color=GRAY)
arrow(s, Inches(2.35), Y3 + BH, Inches(2.35), Y4, color=ORANGE, width=2.5)
arrow(s, g1.left + g1.width, Y4 + Inches(0.56), g2.left, Y4 + Inches(0.56), color=WHITE, width=2.5)
arrow(s, g2.left + g2.width, Y4 + Inches(0.56), g3.left, Y4 + Inches(0.56), color=WHITE, width=2.5)
notes(s, "Three user-facing surfaces share one intelligence layer and one data layer. The key "
         "architectural decision: every path that could contact the outside world funnels "
         "through the orange human-approval layer at the bottom.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — DATA FOUNDATION & INGESTION
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, 4, "Data Foundation — Official Sources → One Vector Store",
       "scripts/ingest.py parses PDFs, CSVs & Excel into 6 ChromaDB collections")
cols = [
    ("hp_hospitals", "289 NHP hospitals (CSV) — district, type, contact, specialities", TEAL),
    ("hp_shelters", "DAY-NULM urban shelters (PDF) — city, capacity", TEAL),
    ("hp_schools", "Govt schools 2021 (PDF) — backup shelter proxies", TEAL),
    ("hp_cwc_stations", "CWC flood-forecast stations (Excel) — river, coordinates", SKY),
    ("hp_disaster_knowledge", "HIMCOSTE 2023 landslide inventory + NDMA/IMD guidance", SKY),
    ("hp_glacial_lakes", "CWC Sept-2025 GLOF monitoring — extracted from the PDF by build_glacial_lakes_csv.py", SKY),
]
x0, y0 = Inches(0.35), Inches(1.32)
cw, ch = Inches(4.15), Inches(1.12)
for i, (name, desc, col) in enumerate(cols):
    x = x0 + (i % 3) * (cw + Inches(0.10))
    y = y0 + (i // 3) * (ch + Inches(0.12))
    b = rect(s, x, y, cw, ch, fill=WHITE, line=col)
    set_text(b, [
        [(name, {"bold": True, "size": 12, "color": col})],
        [(desc, {"size": 9.5, "color": INK})],
    ], anchor=MSO_ANCHOR.MIDDLE, space_after=2)
card(s, Inches(0.35), Inches(3.85), Inches(6.25), Inches(3.0),
     "Embedding & retrieval choices", [
    "sentence-transformers all-MiniLM-L6-v2 — 384-dim, runs locally, fast enough to embed the full corpus in minutes.",
    "ChromaDB PersistentClient — an embedded vector DB: no server to run, survives restarts in chroma_db/.",
    "Metadata filters (where district = KULLU) give exact district scoping on top of semantic search — with an unfiltered fallback if a district returns nothing.",
    "python scripts/ingest.py --overwrite rebuilds everything reproducibly.",
], accent=TEAL)
card(s, Inches(6.85), Inches(3.85), Inches(6.15), Inches(3.0),
     "Two datasets stay outside the vector store — on purpose", [
    "VIIRS wildfire hotspot history (CSV) → loaded as numpy arrays; proneness is a vectorised haversine density count, not a semantic query.",
    "needs.csv / resources.csv (volunteer pool) → pandas records; allocation must be exact-match arithmetic, not fuzzy retrieval.",
    [("Principle: ", {"bold": True, "size": 10.5, "color": NAVY}),
     ("embeddings where language matters, plain math where numbers matter.",
      {"size": 10.5, "color": INK, "italic": True})],
], accent=AMBER)
notes(s, "Emphasise the deliberate split: semantic retrieval for descriptive records, numpy/pandas "
         "for spatial and quantitative data. The GLOF CSV is extracted from CWC's own monitoring "
         "PDF and matches their reported figure exactly (12 expanding HP lakes).")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — LANGGRAPH WORKFLOW
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, 5, "The LangGraph Workflow — a 6-Node State Machine",
       "workflow.py · every node reads & writes one shared DisasterState (TypedDict, ~30 fields)")
nodes = [
    ("1 · intake_agent", "weather + IMD alert · district risk · wildfire flag"),
    ("2 · glof_monitor", "expanding glacial lakes → WATCH / ADVISORY"),
    ("3 · resource_finder", "ChromaDB: hospitals · shelters · CWC stations"),
    ("4 · matching_agent", "LLM ranks resources → priority pick (JSON)"),
    ("5 · route_planner", "ORS route + NH road-risk warnings"),
    ("6 · escalation_agent", "urgency score · final report · contacts"),
]
bx, by = Inches(0.35), Inches(1.45)
bw, bh, gap = Inches(1.98), Inches(1.15), Inches(0.15)
for i, (t, sub) in enumerate(nodes):
    x = bx + i * (bw + gap)
    fill = ORANGE if i in (3, 5) else NAVY
    flow_box(s, x, by, bw, bh, t, sub, fill=fill, title_size=10.5, sub_size=8)
    if i:
        arrow(s, x - gap, by + bh / 2, x, by + bh / 2, color=GRAY, width=2.2)
# conditional skip edge
skip_y = by + bh + Inches(0.32)
arrow(s, bx + 3 * (bw + gap) + bw / 2, by + bh, bx + 3 * (bw + gap) + bw / 2, skip_y, color=RED, width=2, dashed=True)
arrow(s, bx + 3 * (bw + gap) + bw / 2, skip_y, bx + 5 * (bw + gap) + bw / 2, skip_y, color=RED, width=2, dashed=True)
arrow(s, bx + 5 * (bw + gap) + bw / 2, skip_y, bx + 5 * (bw + gap) + bw / 2, by + bh + Inches(0.02), color=RED, width=2, dashed=True)
lb = s.shapes.add_textbox(bx + 3 * (bw + gap) + Inches(0.4), skip_y + Inches(0.02), Inches(4.2), Inches(0.3))
set_text(lb, [[("conditional edge: no resource found → skip routing, escalate",
                {"size": 9.5, "color": RED, "bold": True, "italic": True})]])
card(s, Inches(0.35), Inches(3.5), Inches(6.25), Inches(3.35),
     "Why LangGraph (and not CrewAI / a plain chain)", [
    "Genuine conditional branches: flood → CWC river checks; no match → straight to escalation. Sequential task chaining hides this.",
    "State must persist and accumulate: weather found in node 1 is still needed by the report in node 6 — one typed DisasterState carries it.",
    "add_conditional_edges makes the branch logic explicit, testable code — should_skip_to_escalation() is a 4-line pure function.",
    "A disaster_type_router edge is already wired in — future disaster-specific sub-agents plug in without rewiring.",
    "Human-in-the-loop attaches naturally after the terminal node.",
], accent=ORANGE)
card(s, Inches(6.85), Inches(3.5), Inches(6.15), Inches(3.35),
     "State design details that pay off", [
    "DisasterState is a TypedDict — every field a node fills is declared: inputs, enriched context, resources, matches, route, report.",
    [("node_log: Annotated[list[str], operator.add]", {"size": 10.5, "color": SKY, "bold": True}),
     (" — each node appends a ✓/⚠ line; LangGraph merges them into a built-in agent trace shown in the UI.", {"size": 10.5, "color": INK})],
    "Nodes return {**state, ...updates} — pure functions of state, so each is unit-testable in isolation (18 pytest scenarios).",
    "graph.compile() yields an app; run_agent() seeds the initial state and invokes it end-to-end.",
], accent=SKY)
notes(s, "This slide answers 'why LangGraph' with the exact reasons documented in workflow.py's "
         "docstring. The dashed red edge is the money shot: real branching, not a linear chain.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — NODES I: INTAKE + GLOF
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, 6, "Inside the Nodes I — Situation Intelligence",
       "intake_agent + glof_monitor_agent: turn a location into risk context")
card(s, Inches(0.35), Inches(1.3), Inches(6.25), Inches(5.55),
     "intake_agent — enrich the report", [
    [("Weather → alert:", {"bold": True, "size": 11, "color": NAVY}),
     ("  Open-Meteo 24h precipitation is mapped to IMD-style levels:", {"size": 10.5, "color": INK})],
    [("      RED ≥ 204.5 mm · ORANGE ≥ 115.6 · YELLOW ≥ 64.5 · GREEN below",
      {"size": 10.5, "color": RED, "bold": True})],
    [("District risk:", {"bold": True, "size": 11, "color": NAVY}),
     ("  HIMCOSTE 2023 landslide counts → CRITICAL/HIGH/MEDIUM/LOW tier per district, with key rivers.", {"size": 10.5, "color": INK})],
    [("Nearest CWC station:", {"bold": True, "size": 11, "color": NAVY}),
     ("  haversine scan over all river-monitoring stations → link to live flood forecasts.", {"size": 10.5, "color": INK})],
    [("Wildfire proneness:", {"bold": True, "size": 11, "color": NAVY}),
     ("  vectorised numpy haversine over every historical VIIRS fire hotspot; count within 10 km →", {"size": 10.5, "color": INK})],
    [("      HIGH ≥ 100 hotspots · MODERATE ≥ 25 · LOW ≥ 1 · MINIMAL 0",
      {"size": 10.5, "color": ORANGE, "bold": True})],
    [("Knowledge chunks:", {"bold": True, "size": 11, "color": NAVY}),
     ("  RAG query for district + disaster-type guidance feeds later prompts.", {"size": 10.5, "color": INK})],
    [("Resilience: if the weather API fails, a clearly-labelled MOCK fallback keeps the pipeline alive.",
      {"size": 10, "color": GRAY, "italic": True})],
], accent=SKY, body_size=10.5)
card(s, Inches(6.85), Inches(1.3), Inches(6.15), Inches(5.55),
     "glof_monitor_agent — early warning for a Himalayan killer", [
    "Glacial Lake Outburst Floods gave HP some of its worst recent disasters — so GLOF gets a dedicated node.",
    "Queries the hp_glacial_lakes collection: lakes ranked by (same district, expanding status, proximity).",
    [("Any lake with an increasing water-spread area triggers an alert:", {"size": 10.5, "color": INK})],
    [("      WATCH", {"bold": True, "size": 11, "color": RED}),
     ("  if the disaster is water-driven (GLOF, Flash Flood, Cloudburst)", {"size": 10.5, "color": INK})],
    [("      ADVISORY", {"bold": True, "size": 11, "color": AMBER}),
     ("  otherwise — the hazard exists but isn't the active event", {"size": 10.5, "color": INK})],
    "The alert message names the nearest expanding lake, its basin/river, distance, and % area change — e.g. “+9% water-spread, ~14 km away, elevated GLOF risk downstream”.",
    "Honesty baked in: every alert carries a disclaimer that it is based on CWC's previous-year monthly satellite monitoring (Sep 2025), not real-time water levels.",
    "Validation: the extraction pipeline found 12 expanding HP lakes — exactly CWC's own published figure.",
], accent=RED, body_size=10.5)
notes(s, "Two nodes, one job: convert (lat, lon, district, disaster type) into calibrated risk "
         "signals. Every threshold shown here is real code, not hand-waving.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — NODES II: RESOURCE FINDER + MATCHING
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, 7, "Inside the Nodes II — Finding & Ranking Resources",
       "resource_finder_agent + matching_agent: retrieval is conditional, ranking is LLM with a safety net")
card(s, Inches(0.35), Inches(1.3), Inches(6.25), Inches(5.55),
     "resource_finder_agent — retrieval driven by need & disaster", [
    [("Medical need or injury-prone disaster", {"bold": True, "size": 10.5, "color": NAVY}),
     ("  →  top-5 hospitals (district-filtered semantic query)", {"size": 10.5, "color": INK})],
    [("Shelter need or displacement disaster", {"bold": True, "size": 10.5, "color": NAVY}),
     ("  →  NULM shelters + govt schools as backup shelters", {"size": 10.5, "color": INK})],
    [("Flash Flood / GLOF / Cloudburst", {"bold": True, "size": 10.5, "color": NAVY}),
     ("  →  CWC stations on the nearest river for live-level lookup", {"size": 10.5, "color": INK})],
    "Each query combines semantic text (“hospital emergency KULLU…”) with a hard metadata filter on district — and falls back to an unfiltered query if the district yields nothing.",
    "Also pulls 3 extra knowledge chunks (“<disaster> response resources <district>”) and merges them into state for the prompts downstream.",
    [("Node log: ", {"bold": True, "size": 10, "color": GRAY}),
     ("“✓ found 5 hospitals, 3 shelters, 2 CWC stations” — the trace users see.",
      {"size": 10, "color": GRAY, "italic": True})],
], accent=TEAL, body_size=10.5)
card(s, Inches(6.85), Inches(1.3), Inches(6.15), Inches(5.55),
     "matching_agent — LLM ranking with a deterministic parachute", [
    "All retrieved hospitals + shelters are pooled and tagged with resource_type.",
    [("Empty pool? ", {"bold": True, "size": 10.5, "color": RED}),
     ("Set escalation_needed=True and return — the conditional edge then skips routing entirely.", {"size": 10.5, "color": INK})],
    "Otherwise llama3.2:1b gets the situation (district, disaster, needs, IMD alert, risk tier), the full resource JSON, and knowledge chunks…",
    [("…and must answer in a strict JSON contract: ", {"size": 10.5, "color": INK}),
     ("priority_resource · ranked_resources · 2-sentence reasoning", {"size": 10.5, "color": SKY, "bold": True})],
    "Temperature 0.1 keeps it near-deterministic; the reply is stripped of ``` fences and json.loads-ed.",
    [("If parsing fails: ", {"bold": True, "size": 10.5, "color": ORANGE}),
     ("fall back to retrieval order — the top semantic match becomes priority, and the reasoning says a fallback was used. The pipeline never crashes on a bad LLM reply.", {"size": 10.5, "color": INK})],
    [("Pattern: LLM for judgement, code for guarantees.",
      {"bold": True, "italic": True, "size": 11, "color": NAVY})],
], accent=ORANGE, body_size=10.5)
notes(s, "The matching agent is where the LLM adds judgement — weighing an ORANGE alert and a "
         "CRITICAL district against resource types. But the JSON contract + fallback means a "
         "1B-parameter local model can never take the pipeline down.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — NODES III: ROUTE + ESCALATION / URGENCY
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, 8, "Inside the Nodes III — Routes, Urgency & the Final Report",
       "route_planning_agent + escalation_agent: always produce a usable answer")
card(s, Inches(0.35), Inches(1.3), Inches(6.25), Inches(2.9),
     "route_planning_agent — a destination ladder that never fails", [
    [("Resolve where the resource actually is, best-effort:", {"bold": True, "size": 10.5, "color": NAVY})],
    [("  resource coords → geocode full name → geocode the town inside the name "
      "(“Civil Hospital Dalhousie” → “Dalhousie”) → district centre",
      {"size": 10.5, "color": SKY, "bold": True})],
    "Routing: OpenRouteService turn-by-turn; if ORS is down/keyless → straight-line haversine at ~30 km/h mountain speed, explicitly labelled approximate.",
    "Computes routes to the top hospital AND top shelter, plus NH corridor warnings (Kotrupi, Nigulsari…) during monsoon months.",
], accent=SKY, body_size=10.5)
card(s, Inches(0.35), Inches(4.35), Inches(6.25), Inches(2.5),
     "escalation_agent — the report the human reads", [
    "llama3.2:1b writes a markdown situation report from the full state: weather, risks, priority resource, route, GLOF/wildfire flags, top-3 alternatives.",
    "Always appends HP emergency contacts (NDMA 1078, Police 100, 108, NDRF…) and the nearest hospital's phone number.",
    "Every report ends: “AI-assisted information — always verify with local authorities.”",
], accent=ORANGE, body_size=10.5)
# Urgency panel
up = rect(s, Inches(6.85), Inches(1.3), Inches(6.15), Inches(5.55), fill=NAVY)
set_text(up, [
    [("Explainable urgency score (0–100) — pure arithmetic, no LLM", {"bold": True, "size": 14, "color": ORANGE})],
    [("", {"size": 4})],
    [("IMD alert", {"bold": True, "size": 12, "color": WHITE}),
     ("   RED 35 · ORANGE 25 · YELLOW 15 · GREEN 5", {"size": 11.5, "color": RGBColor(0xC9, 0xD8, 0xE8)})],
    [("District tier", {"bold": True, "size": 12, "color": WHITE}),
     ("   CRITICAL 25 · HIGH 18 · MEDIUM 10 · LOW 5", {"size": 11.5, "color": RGBColor(0xC9, 0xD8, 0xE8)})],
    [("Needs (top 2)", {"bold": True, "size": 12, "color": WHITE}),
     ("   Rescue/Medical 10 · Evacuation 9 · Shelter/Water 6 · Food 4", {"size": 11.5, "color": RGBColor(0xC9, 0xD8, 0xE8)})],
    [("GLOF signal", {"bold": True, "size": 12, "color": WHITE}),
     ("   WATCH 12 · ADVISORY 6", {"size": 11.5, "color": RGBColor(0xC9, 0xD8, 0xE8)})],
    [("Wildfire", {"bold": True, "size": 12, "color": WHITE}),
     ("   HIGH 10 · MODERATE 6 · LOW 2", {"size": 11.5, "color": RGBColor(0xC9, 0xD8, 0xE8)})],
    [("Escalation flag", {"bold": True, "size": 12, "color": WHITE}),
     ("   +10 when no resource could be matched", {"size": 11.5, "color": RGBColor(0xC9, 0xD8, 0xE8)})],
    [("", {"size": 4})],
    [("Σ capped at 100  →  CRITICAL ≥ 75 · HIGH ≥ 50 · MODERATE ≥ 30 · LOW",
      {"bold": True, "size": 12.5, "color": AMBER})],
    [("", {"size": 4})],
    [("The per-factor breakdown ships with the score, so a coordinator can see exactly why "
      "one incident outranks another — and challenge it.",
      {"size": 11, "color": WHITE, "italic": True})],
], anchor=MSO_ANCHOR.TOP, space_after=7)
notes(s, "Two deliberate design choices: the routing ladder guarantees a number always appears "
         "(clearly labelled when approximate), and urgency is deterministic arithmetic so "
         "prioritisation between incidents is reproducible and defensible.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — RESOURCE ALLOCATION ENGINE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, 9, "Allocating Resources to People — the Matching Engine",
       "coordination.py · every reported need is scored against every available provider — deterministically")
# scoring bar visual
sb_y = Inches(1.35)
segs = [("Category match  40", Inches(4.6), NAVY),
        ("Location overlap  ≤25", Inches(2.9), SKY),
        ("Quantity coverage  ≤20", Inches(2.3), TEAL),
        ("Available  8", Inches(1.15), AMBER),
        ("Verified  7", Inches(1.05), ORANGE)]
sx = Inches(0.65)
for label, w, col in segs:
    b = rect(s, sx, sb_y, w, Inches(0.62), fill=col, shape=MSO_SHAPE.RECTANGLE)
    set_text(b, [[(label, {"bold": True, "size": 10.5, "color": WHITE})]], align=PP_ALIGN.CENTER)
    sx += w
tb = s.shapes.add_textbox(Inches(0.65), sb_y + Inches(0.66), Inches(12.0), Inches(0.35))
set_text(tb, [[("100-point match score — every point traceable to a rule, so two coordinators always get the same allocation",
                {"size": 10.5, "color": GRAY, "italic": True})]])
card(s, Inches(0.35), Inches(2.55), Inches(6.25), Inches(4.3),
     "How a request becomes an allocation", [
    [("1.  ", {"bold": True, "color": ORANGE, "size": 10.5}),
     ("Requests arrive from needs.csv or free text/tweets — a rule-based extractor pulls category, urgency, quantity and location from messages like “30 people trapped near Aut, need rescue urgently”.", {"size": 10.5, "color": INK})],
    [("2.  ", {"bold": True, "color": ORANGE, "size": 10.5}),
     ("Category is a hard filter — a Food provider can never be allocated to a Medical need.", {"size": 10.5, "color": INK})],
    [("3.  ", {"bold": True, "color": ORANGE, "size": 10.5}),
     ("Location score = token overlap of place names (Jaccard) · quantity score = min(available/needed, 1).", {"size": 10.5, "color": INK})],
    [("4.  ", {"bold": True, "color": ORANGE, "size": 10.5}),
     ("Ties break on provider urgency-capacity; each need keeps its best match + 2 alternatives.", {"size": 10.5, "color": INK})],
    [("5.  ", {"bold": True, "color": ORANGE, "size": 10.5}),
     ("The worklist is sorted by need urgency (Critical first), unmatched needs surfaced first within each level — triage order, not file order.", {"size": 10.5, "color": INK})],
], accent=ORANGE, body_size=10.5)
card(s, Inches(6.85), Inches(2.55), Inches(6.15), Inches(4.3),
     "Allocation outcomes", [
    [("MATCHED", {"bold": True, "size": 11.5, "color": TEAL}),
     ("  score ≥ 55 — draft outreach to the provider for min(needed, available) units.", {"size": 10.5, "color": INK})],
    [("PARTIAL", {"bold": True, "size": 11.5, "color": AMBER}),
     ("  best match exists but is weak/short — coverage % and the exact unit shortfall are shown, with “consider a second provider”.", {"size": 10.5, "color": INK})],
    [("UNMATCHED", {"bold": True, "size": 11.5, "color": RED}),
     ("  nothing eligible — an escalation draft to HPSDMA 1077 / NDMA 1078 is generated instead. No request silently disappears.", {"size": 10.5, "color": INK})],
    [("", {"size": 3})],
    [("Why no LLM here: ", {"bold": True, "size": 11, "color": NAVY}),
     ("allocation decides who gets scarce resources. It must be transparent, reproducible and auditable — properties arithmetic has and generative models don't. The LLM assists around the decision; it never makes it.",
      {"size": 10.5, "color": INK, "italic": True})],
    [("Validated: all 5 sample needs matched to the correct provider in tests.",
      {"size": 10, "color": GRAY, "italic": True})],
], accent=TEAL, body_size=10.5)
notes(s, "This is the heart of 'how resources are allocated based on a request'. Walk the scoring "
         "bar left to right, then the three outcome states. Land the punchline: LLM for language, "
         "arithmetic for allocation.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — MATCHING, WORKED EXAMPLE (one real need, fully scored)
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, 10, "Matching, Step by Step — One Need, Fully Scored",
       "coordination.py · a real request from needs.csv scored against every eligible provider")
# ── Incoming need strip ────────────────────────────────────────────────
need_strip = rect(s, Inches(0.35), Inches(1.24), Inches(12.65), Inches(0.86), fill=NAVY)
set_text(need_strip, [
    [("Incoming need  N003", {"bold": True, "size": 12.5, "color": ORANGE}),
     ("     Category ", {"size": 11.5, "color": WHITE}), ("Food", {"bold": True, "size": 11.5, "color": WHITE}),
     ("      Location ", {"size": 11.5, "color": WHITE}), ("Mandi riverbank", {"bold": True, "size": 11.5, "color": WHITE}),
     ("      Quantity ", {"size": 11.5, "color": WHITE}), ("60 meal kits", {"bold": True, "size": 11.5, "color": WHITE}),
     ("      Reported by ", {"size": 11.5, "color": WHITE}), ("NGO field lead", {"bold": True, "size": 11.5, "color": WHITE})],
    [("Same request can arrive from a tweet — the rule-based extractor pulls category / quantity / location before scoring.",
      {"size": 10, "color": RGBColor(0xC9, 0xD8, 0xE8), "italic": True})],
], anchor=MSO_ANCHOR.MIDDLE, space_after=3)
# ── Scoreboard ─────────────────────────────────────────────────────────
lbl = s.shapes.add_textbox(Inches(0.35), Inches(2.24), Inches(12.6), Inches(0.3))
set_text(lbl, [[("Score every eligible provider — category is a hard filter, so only Food providers are scored:",
                 {"size": 11, "color": GRAY, "italic": True})]], anchor=MSO_ANCHOR.MIDDLE)
sb_cols = [Inches(0.35), Inches(4.15), Inches(5.55), Inches(7.15), Inches(8.75), Inches(10.55)]
sb_w    = [Inches(3.70), Inches(1.30), Inches(1.50), Inches(1.50), Inches(1.70), Inches(2.45)]
def sb_row(y, cells, fill, text_color=INK, bold=False):
    for i, (cx, cw, val) in enumerate(zip(sb_cols, sb_w, cells)):
        c = rect(s, cx, y, cw, Inches(0.52), fill=fill, shape=MSO_SHAPE.RECTANGLE)
        set_text(c, [[(val, {"bold": bold, "size": 10.5, "color": text_color})]],
                 align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
sb_row(Inches(2.60), ["Eligible Food provider", "Category\n40", "Location\n≤25", "Quantity\n≤20", "Avail+Verif\n15", "Score /100"],
       NAVY, text_color=WHITE, bold=True)
sb_row(Inches(3.16), ["Red Cross Mandi  (Mandi)", "40", "12.5", "20", "15", "87.5  ✔ BEST"], AMBER, text_color=NAVY, bold=True)
sb_row(Inches(3.72), ["Annapurna Kitchen  (Shimla)", "40", "0", "20", "15", "75.0"], WHITE)
sb_row(Inches(4.28), ["Solan Wholesale Foods  (Solan)", "40", "0", "20", "15", "75.0"], CARD)
# ── Bottom: why it wins + outcome ──────────────────────────────────────
card(s, Inches(0.35), Inches(5.05), Inches(6.25), Inches(1.85),
     "Why Red Cross Mandi wins", [
    [("Category 40", {"bold": True, "size": 10.5, "color": NAVY}),
     ("  hard filter passed (Food = Food).   ", {"size": 10.5, "color": INK}),
     ("Location 12.5", {"bold": True, "size": 10.5, "color": SKY}),
     ("  “Mandi” token overlaps (0.5 Jaccard × 25).", {"size": 10.5, "color": INK})],
    [("Quantity 20", {"bold": True, "size": 10.5, "color": TEAL}),
     ("  120 kits ≥ 60 needed → full coverage.   ", {"size": 10.5, "color": INK}),
     ("Avail 15", {"bold": True, "size": 10.5, "color": ORANGE}),
     ("  Available (8) + Verified (7).", {"size": 10.5, "color": INK})],
    "The two Shimla/Solan providers tie at 75.0 — identical except location. Proximity is the whole difference: the same rubric, applied identically, breaks the tie.",
], accent=AMBER, body_size=10.5)
card(s, Inches(6.85), Inches(5.05), Inches(6.15), Inches(1.85),
     "Outcome → an actionable draft", [
    [("Score 87.5 ≥ 55  →  ", {"bold": True, "size": 11, "color": TEAL}),
     ("MATCHED", {"bold": True, "size": 11.5, "color": TEAL}),
     (". Draft outreach for min(60, 120) = 60 kits, coverage 100%.", {"size": 10.5, "color": INK})],
    "Best match + 2 ranked alternatives are kept, so a coordinator can override the pick with one click.",
    [("Every number above is reproducible: ", {"italic": True, "size": 10, "color": GRAY}),
     ("no LLM touches the allocation score — two coordinators always see the same 87.5.",
      {"italic": True, "size": 10, "color": GRAY})],
], accent=TEAL, body_size=10.5)
notes(s, "Walk the scoreboard top row (the rubric) then the three rows. The teaching point: the two "
         "alternatives are identical to the winner on everything except location — so proximity alone "
         "decides it. Every value here is real output from coordination.match_needs_to_resources().")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — INVENTORY-AWARE ALLOCATION (no unit promised twice)
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, 11, "Inventory-Aware Allocation — No Unit Promised Twice",
       "coordination.py allocate=True · matching is stateful: process → reserve → re-score against what's left")
card(s, Inches(0.35), Inches(1.3), Inches(3.75), Inches(5.55),
     "Why per-need scoring isn't enough", [
    "Scoring each need independently makes one provider with 100 units look fully available to every need at once.",
    "In a real dispatch that double-promises scarce stock — two coordinators act on the same 100 units.",
    [("Fix: ", {"bold": True, "size": 10.5, "color": ORANGE}),
     ("needs are processed in urgency order, and each best match reserves its units from a working copy of the pool before the next need is scored.",
      {"size": 10.5, "color": INK})],
    "So allocation reflects triage priority AND real remaining stock — not file order, not listed capacity.",
], accent=RED, body_size=10.5)
# ── Center worked ledger ───────────────────────────────────────────────
panel = rect(s, Inches(4.30), Inches(1.3), Inches(4.55), Inches(5.55), fill=NAVY)
set_text(panel, [
    [("One provider · 100 units · 3 competing needs", {"bold": True, "size": 12.5, "color": ORANGE})],
    [("Served in urgency order, not file order:", {"size": 10.5, "color": RGBColor(0xC9, 0xD8, 0xE8), "italic": True})],
    [("", {"size": 4})],
    [("① CRITICAL · needs 80", {"bold": True, "size": 11.5, "color": WHITE})],
    [("     reserve 80  →  20 left   ", {"size": 10.5, "color": RGBColor(0xC9, 0xD8, 0xE8)}),
     ("MATCHED", {"bold": True, "size": 10.5, "color": TEAL})],
    [("", {"size": 3})],
    [("② HIGH · needs 50", {"bold": True, "size": 11.5, "color": WHITE})],
    [("     only 20 left → reserve 20   ", {"size": 10.5, "color": RGBColor(0xC9, 0xD8, 0xE8)}),
     ("PARTIAL", {"bold": True, "size": 10.5, "color": AMBER})],
    [("     shortfall −30 → “consider a second provider”", {"size": 9.5, "color": RGBColor(0xC9, 0xD8, 0xE8)})],
    [("", {"size": 3})],
    [("③ MEDIUM · needs 15", {"bold": True, "size": 11.5, "color": WHITE})],
    [("     provider exhausted, skipped   ", {"size": 10.5, "color": RGBColor(0xC9, 0xD8, 0xE8)}),
     ("UNMATCHED", {"bold": True, "size": 10.5, "color": RED})],
    [("     → escalation draft to HPSDMA 1077 / NDMA 1078", {"size": 9.5, "color": RGBColor(0xC9, 0xD8, 0xE8)})],
    [("", {"size": 5})],
    [("committed_units never exceed what's actually available — total dispatched ≤ 100, always.",
      {"size": 10, "color": AMBER, "italic": True})],
], anchor=MSO_ANCHOR.TOP, space_after=4)
card(s, Inches(9.05), Inches(1.3), Inches(3.95), Inches(5.55),
     "Dispatch ledger — stock stays spent", [
    [("Approve & send → ", {"size": 10.5, "color": INK}),
     ("log_dispatch()", {"bold": True, "size": 10.5, "color": SKY}),
     (" appends {resource_id, units, request_id} to logs/dispatch_ledger.jsonl.", {"size": 10.5, "color": INK})],
    [("load_resources()", {"bold": True, "size": 10.5, "color": SKY}),
     (" subtracts dispatched totals: remaining = listed − dispatched, floored at 0.", {"size": 10.5, "color": INK})],
    "The UI shows listed / dispatched / remaining per provider, so depletion is visible across sessions.",
    "Approvals in one incident correctly shrink the pool for the next — later needs degrade to PARTIAL/UNMATCHED on their own.",
    [("Proven in tests: ", {"bold": True, "size": 10, "color": GRAY, "italic": True}),
     ("no double-promising, urgency-first order, exhaustion, and stock never goes negative.",
      {"size": 10, "color": GRAY, "italic": True})],
], accent=TEAL, body_size=10.5)
notes(s, "This is the second half of 'how matching is done': not just scoring a pair, but allocating a "
         "finite pool across competing needs. Walk the center ledger top to bottom — the same 100 units "
         "produce a MATCHED, a PARTIAL and an UNMATCHED because stock depletes as it's reserved. The "
         "dispatch ledger makes that depletion persist across runs.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 12 — HUMAN-IN-THE-LOOP EMAIL WORKFLOW
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, 12, "From Match to Mailbox — the Human-in-the-Loop Email Flow",
       "Every coordination message is drafted by the system, decided by a person, and logged forever")
steps = [
    ("1 · Draft", "draft_coordination_message() — provider, units, coverage %, confidence /100", NAVY),
    ("2 · Review", "coordinator reads the draft in Streamlit — full text, editable", NAVY),
    ("3 · Decide", "✅ approve & send · ✔ approve only · ❌ reject", ORANGE),
    ("4 · Send", "smtplib → smtp.gmail.com:587 · STARTTLS · Gmail app password", TEAL),
    ("5 · Audit", "decision appended to logs/approvals.jsonl (UTC timestamp)", SKY),
]
bx, by = Inches(0.35), Inches(1.45)
bw, bh, gap = Inches(2.4), Inches(1.3), Inches(0.16)
for i, (t, sub, col) in enumerate(steps):
    x = bx + i * (bw + gap)
    flow_box(s, x, by, bw, bh, t, sub, fill=col, title_size=12, sub_size=8.5)
    if i:
        arrow(s, x - gap, by + bh / 2, x, by + bh / 2, color=GRAY, width=2.2)
card(s, Inches(0.35), Inches(3.2), Inches(6.25), Inches(3.65),
     "What's in the email", [
    "Header: “DISASTER RELIEF COORDINATION — MEDICAL [CRITICAL priority]”.",
    "The need: request ID, quantity, location, who reported it, notes.",
    "The match: provider name & location, listed stock, availability + contact-verification status, estimated coverage %.",
    "The ask: dispatch min(needed, available) units + expected arrival time.",
    "The caveat: “Match confidence 87/100 — requires human coordinator approval before sending.”",
    "Unmatched needs get a different draft: an escalation to the district control room instead of provider outreach.",
], accent=NAVY, body_size=10.5)
card(s, Inches(6.85), Inches(3.2), Inches(6.15), Inches(3.65),
     "Why this design", [
    [("Nothing sends itself. ", {"bold": True, "size": 11, "color": RED}),
     ("Drafts are inert text until a person clicks approve — in a disaster, a wrong automated message costs trust and lives.", {"size": 10.5, "color": INK})],
    "Editable before send — the human can correct quantities or add local knowledge the system lacks.",
    "Approve, reject AND send-failures are all logged to approvals.jsonl — a replayable audit trail of every human decision.",
    "Plain smtplib + EmailMessage: zero vendor lock-in, works with any Gmail app password entered in the sidebar (never stored in code).",
    "The same gate wraps both surfaces: volunteer-worklist matches and the LangGraph agent's own coordination draft.",
], accent=ORANGE, body_size=10.5)
notes(s, "Demo beat: click approve on a match, show the real email arriving, then open "
         "approvals.jsonl to show the logged decision. Rejections are logged too — the audit "
         "trail captures what the human declined, not just what they sent.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 13 — TECH STACK: WHY EACH PIECE
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, 13, "The Stack — Every Choice Has a Reason",
       "Local-first, free-tier, zero mandatory API keys — built to run in a district control room")
rows = [
    ("LangGraph", "Explicit state machine: typed shared state, conditional edges, natural human-in-the-loop hook — branching CrewAI-style chains can't express."),
    ("LangChain (core + ollama)", "Uniform chat-model interface & message types; swap the LLM without touching agent code."),
    ("Ollama · llama3.2:1b @ t=0.1", "Fully local & free — no cloud dependency when connectivity is down (exactly when disasters strike); no sensitive location data leaves the machine."),
    ("ChromaDB", "Embedded, persistent vector DB — no server to operate; metadata filters give exact district scoping on top of semantic search."),
    ("sentence-transformers (MiniLM-L6-v2)", "384-dim local embeddings: fast, offline, strong enough that the RAG relevance gate separates in/out-of-scope by a 0.18 distance margin."),
    ("Streamlit", "A control-room UI in pure Python; session_state drives the approve/reject workflow without a JS frontend."),
    ("pandas + numpy", "CSV/Excel wrangling; vectorised haversine scores thousands of VIIRS fire hotspots in milliseconds."),
    ("folium", "Interactive Leaflet maps of user, resources and routes — embedded straight into Streamlit."),
    ("Open-Meteo", "Free, key-less forecast API → 24h rainfall → IMD-style alert derivation; mock fallback keeps demos alive offline."),
    ("OpenRouteService + Nominatim", "Free-tier road routing and geocoding; both wrapped with labelled fallbacks (straight-line estimate, district centre)."),
    ("smtplib / EmailMessage (stdlib)", "Email dispatch with no third-party service — one less credential, one less point of failure."),
    ("pytest + loguru", "18 deterministic scenario tests pin urgency bands, matching, and the RAG gate; structured logs for every tool failure."),
]
ty = Inches(1.28)
th = Inches(0.462)
for i, (tech, why) in enumerate(rows):
    y = ty + i * th
    band_fill = WHITE if i % 2 == 0 else CARD
    r1 = rect(s, Inches(0.35), y, Inches(3.35), th, fill=band_fill, shape=MSO_SHAPE.RECTANGLE)
    r2 = rect(s, Inches(3.70), y, Inches(9.28), th, fill=band_fill, shape=MSO_SHAPE.RECTANGLE)
    set_text(r1, [[(tech, {"bold": True, "size": 10.5, "color": NAVY})]])
    set_text(r2, [[(why, {"size": 9.5, "color": INK})]])
tb = s.shapes.add_textbox(Inches(0.35), ty + 12 * th + Inches(0.02), Inches(12.6), Inches(0.28))
set_text(tb, [[("Common thread: ", {"bold": True, "size": 10.5, "color": ORANGE}),
               ("everything critical runs locally and degrades gracefully — the agent still answers when the internet doesn't.",
                {"size": 10.5, "color": INK, "italic": True})]], anchor=MSO_ANCHOR.TOP)
notes(s, "If asked 'why not GPT-4/cloud': disasters knock out connectivity, and location reports "
         "are sensitive. A 1B local model is enough because the LLM only ranks and writes — all "
         "safety-critical logic is deterministic code.")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 14 — GUARDRAILS, VALIDATION, ROADMAP
# ══════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, 14, "Trust, Proof & What's Next",
       "Guardrails you can measure, results you can reproduce")
card(s, Inches(0.35), Inches(1.3), Inches(4.15), Inches(4.1),
     "Guardrails", [
    "Human approval gates every outbound message; full decision audit in approvals.jsonl.",
    "RAG chatbot refuses out-of-scope questions before the LLM runs (relevance gate at 0.70 cosine distance) and is prompt-locked to retrieved context only.",
    "Every fallback is labelled: MOCK weather, straight-line routes, previous-year GLOF/wildfire data disclaimers.",
    "Every report ends with “verify with local authorities”.",
], accent=RED, body_size=10)
card(s, Inches(4.60), Inches(1.3), Inches(4.15), Inches(4.1),
     "Measured results", [
    [("18/18", {"bold": True, "size": 13, "color": TEAL}),
     ("  deterministic pytest scenarios pass — urgency bands, matching, escalation, RAG gate.", {"size": 10, "color": INK})],
    [("0", {"bold": True, "size": 13, "color": TEAL}),
     ("  false accepts/rejects at the RAG gate: in-scope ≤ 0.59 vs out-of-scope ≥ 0.77 distance.", {"size": 10, "color": INK})],
    [("12", {"bold": True, "size": 13, "color": TEAL}),
     ("  expanding glacial lakes extracted — exactly matches CWC's own published count.", {"size": 10, "color": INK})],
    [("5/5", {"bold": True, "size": 13, "color": TEAL}),
     ("  sample needs allocated to the correct provider.", {"size": 10, "color": INK})],
    [("node_log", {"bold": True, "size": 11, "color": SKY}),
     ("  gives a step-by-step agent trace for every run.", {"size": 10, "color": INK})],
], accent=TEAL, body_size=10)
card(s, Inches(8.85), Inches(1.3), Inches(4.15), Inches(4.1),
     "Roadmap", [
    "Live CWC river levels & IMD bulletins instead of derived alerts.",
    "Verified resource coordinates → true door-to-door routing.",
    "SMS / WhatsApp dispatch behind the same approval gate.",
    "Disaster-specific sub-agents on the already-wired disaster_type_router edge.",
    "Real-time volunteer availability sync.",
], accent=AMBER, body_size=10)
band = rect(s, Inches(0.35), Inches(5.6), Inches(12.65), Inches(1.3), fill=DARKNAVY)
set_text(band, [
    [("ResQ", {"bold": True, "size": 16, "color": ORANGE}),
     ("  ·  Retrieval finds the options. Deterministic scoring allocates them. A local LLM explains them. A human approves them.",
      {"size": 14, "color": WHITE})],
    [("AI does the triage math in seconds — people keep the authority. Thank you.",
      {"size": 12, "color": RGBColor(0xBF, 0xD3, 0xE6), "italic": True})],
], anchor=MSO_ANCHOR.MIDDLE, space_after=4)
notes(s, "Close on the division of labour sentence — it summarises the whole architecture in one "
         "line and answers the 'why is this responsible AI' question before it's asked.")

# ── Save ───────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parents[1] / "docs" / "ResQ_Capstone_Presentation.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"Saved {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
